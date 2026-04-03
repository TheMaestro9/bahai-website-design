import pptx
import sys
import os

# Set stdout to use utf-8 to avoid issues with non-ASCII characters
if sys.version_info >= (3, 7):
    sys.stdout.reconfigure(encoding='utf-8')

def extract_text(file_path):
    if not os.path.exists(file_path):
        print(f"Error: File '{file_path}' not found.")
        return

    try:
        prs = pptx.Presentation(file_path)
    except Exception as e:
        print(f"Error loading PPTX: {e}")
        return

    print(f"### Presentation: {os.path.basename(file_path)}")
    print(f"Number of slides: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        try:
            # Extract title
            if slide.shapes.title:
                print(f"Title: {slide.shapes.title.text.strip()}\n")
            
            # Extract all text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if slide.shapes.title and shape == slide.shapes.title:
                        continue
                    text = shape.text.strip()
                    # Clean up common PPTX artifacts
                    text = text.replace('\r', '\n').strip()
                    print(text)
            
            # Extract notes if any
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    print(f"\nNotes: {notes}")
        except Exception as e:
            print(f"Exception on Slide {i+1}: {e}")
        print("\n")

if __name__ == "__main__":
    pptx_file = "Home Page -team.pptx"
    extract_text(pptx_file)

